import os
from groq import Groq
from pypdf import PdfReader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from pptx import Presentation
from docx import Document as DocxDocument
VECTORSTORE = None
DB_PATH = "vectorstore"

def carregar_memoria():
    global VECTORSTORE

    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    if os.path.exists(DB_PATH):
        VECTORSTORE = FAISS.load_local(DB_PATH, embeddings, allow_dangerous_deserialization=True) 
def ler_arquivo(file):
    docs = []

    # PDF
    if file.type == "application/pdf":
        reader = PdfReader(file)

        for page in reader.pages:
            texto = page.extract_text()
            if texto:
                texto = " ".join(texto.split())
                docs.append(Document(page_content=texto))

    # TXT
    elif file.type == "text/plain":
        texto = str(file.read(), "utf-8")
        docs.append(Document(page_content=texto))

    # DOCX
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = DocxDocument(file)

        texto = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        docs.append(Document(page_content=texto))

    # PPTX
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(file)

        texto_total = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_total += shape.text + "\n"

        docs.append(Document(page_content=texto_total))

    return docs
    

def quebrar_texto(docs):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=900,
        chunk_overlap=200
    )
    return splitter.split_documents(docs)


def filtrar_chunks(chunks):
    palavras_ruins = [
        "introdução",
        "sumário",
        "agradecimento",
        "manual de operação",
        "anúncio"
    ]

    filtrados = []

    for c in chunks:
        texto = c.page_content.lower()
        if not any(p in texto for p in palavras_ruins):
            filtrados.append(c)

    return filtrados


def salvar_memoria(chunks):
    global VECTORSTORE

    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    chunks = filtrar_chunks(chunks)
    VECTORSTORE = FAISS.from_documents(chunks, embeddings)
    
    VECTORSTORE.save_local(DB_PATH)
    
def buscar_memoria(pergunta):
    global VECTORSTORE

    if VECTORSTORE is None:
        return "Nenhum documento carregado."

    pergunta = pergunta.lower().strip()

    stopwords = ["qual", "quais", "o", "a", "os", "as", "de", "da", "do", "das", "dos", "é", "em", "para", "com"]

    palavras = [
        p for p in pergunta.split()
        if p not in stopwords and len(p) > 2
    ]

    # busca inicial
    docs = VECTORSTORE.similarity_search(pergunta, k=30)

    # filtro simples
    docs_com_score = []

    for doc in docs:
        texto = doc.page_content.lower()

        score = sum(1 for p in palavras if p in texto)

        if score > 0:
            docs_com_score.append((doc, score))

    
    docs_com_score = sorted(docs_com_score, key=lambda x: x[1], reverse=True)

    if docs_com_score:
        docs = [d[0] for d in docs_com_score]

    docs = docs[:15]

    contexto = "\n\n".join([
    doc.page_content
    for doc in docs
    ])


    print(contexto)

    client = Groq(
        api_key=os.environ.get("GROQ_API_KEY")
    )

    prompt = f"""
Você é um especialista em manutenção industrial.

INSTRUÇÕES:

- Use APENAS o TEXTO fornecido
- Responda de forma técnica, completa e clara
- Traga TODAS as informações relevantes encontradas
- NÃO limite o tamanho da resposta
- NÃO invente informações

Se a resposta NÃO estiver no texto:
inicie com:
"SUGESTÃO BASEADA EM CONHECIMENTO TÉCNICO DE MERCADO"

PERGUNTA:
{pergunta}

TEXTO:
{contexto}
""" 

    chat = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant"
    )

    resposta = chat.choices[0].message.content

    return resposta
